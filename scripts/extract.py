#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import re
import sys
import shutil
import logging
import argparse
from pathlib import Path
import traceback
from typing import Dict, List, Optional, Tuple

# PDF processing
import pdfplumber
from pdf2image import convert_from_path

# OCR
import pytesseract
from PIL import Image

# Word documents
import docx2txt

# PowerPoint
try:
    from pptx import Presentation
except ImportError:
    pass

# For .doc and .ppt
try:
    import textract
except ImportError:
    pass

# For Windows with Office
try:
    import win32com.client
    OFFICE_INSTALLED = True
except ImportError:
    OFFICE_INSTALLED = False

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler("extraction.log", mode='w')
    ]
)
logger = logging.getLogger()

def setup_argparse() -> argparse.Namespace:
    """Set up command line arguments"""
    parser = argparse.ArgumentParser(description="Extract text from document files")
    parser.add_argument("input_dir", type=str, help="Input directory with documents")
    parser.add_argument("output_dir", type=str, help="Output directory for extracted text")
    return parser.parse_args()

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF, using OCR if needed"""
    try:
        # Try pdfplumber first (for digital PDFs)
        with pdfplumber.open(file_path) as pdf:
            text = ""
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text += page_text + "\n\n"
            
            # If no text was extracted, PDF might be scanned - use OCR
            if not text.strip():
                logger.info(f"PDF appears to be scanned, using OCR: {file_path}")
                images = convert_from_path(file_path)
                for img in images:
                    text += pytesseract.image_to_string(img, lang='fra+eng') + "\n\n"
            
            return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path}: {str(e)}")
        return ""

def extract_text_from_image(file_path: str) -> str:
    """Extract text from image using OCR"""
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang='fra+eng')
        return text
    except Exception as e:
        logger.error(f"Error extracting text from image {file_path}: {str(e)}")
        return ""

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from .docx file"""
    try:
        text = docx2txt.process(file_path)
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {str(e)}")
        return ""

def extract_text_from_doc(file_path: str) -> str:
    """Extract text from .doc file using available methods"""
    try:
        # Try textract if available
        try:
            text = textract.process(file_path).decode('utf-8')
            return text
        except (ImportError, Exception) as e:
            logger.warning(f"Textract failed for {file_path}: {str(e)}")
            
        # Try win32com if Office is installed
        if OFFICE_INSTALLED:
            try:
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(os.path.abspath(file_path))
                text = doc.Content.Text
                doc.Close()
                word.Quit()
                return text
            except Exception as e:
                logger.warning(f"Win32com failed for {file_path}: {str(e)}")
        
        logger.error(f"Failed to extract text from DOC {file_path} with all methods")
        return ""
    except Exception as e:
        logger.error(f"Error extracting text from DOC {file_path}: {str(e)}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from .pptx file"""
    try:
        text = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {str(e)}")
        return ""

def extract_text_from_ppt(file_path: str) -> str:
    """Extract text from .ppt file using available methods"""
    try:
        # Try textract if available
        try:
            text = textract.process(file_path).decode('utf-8')
            return text
        except (ImportError, Exception) as e:
            logger.warning(f"Textract failed for {file_path}: {str(e)}")
            
        # Try win32com if Office is installed
        if OFFICE_INSTALLED:
            try:
                ppt = win32com.client.Dispatch("PowerPoint.Application")
                presentation = ppt.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
                text = ""
                for slide in presentation.Slides:
                    for shape in slide.Shapes:
                        if shape.HasTextFrame:
                            if shape.TextFrame.HasText:
                                text += shape.TextFrame.TextRange.Text + "\n\n"
                presentation.Close()
                ppt.Quit()
                return text
            except Exception as e:
                logger.warning(f"Win32com failed for {file_path}: {str(e)}")
        
        logger.error(f"Failed to extract text from PPT {file_path} with all methods")
        return ""
    except Exception as e:
        logger.error(f"Error extracting text from PPT {file_path}: {str(e)}")
        return ""

def clean_text(text: str) -> str:
    """Clean extracted text"""
    # Replace multiple spaces with a single space
    text = re.sub(r'\s+', ' ', text)
    
    # Normalize line breaks
    text = re.sub(r'\n\s*\n+', '\n\n', text)
    
    # Remove control characters
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
    
    # Remove page numbers (various formats)
    text = re.sub(r'\n\s*\d+\s*\n', '\n', text)  # Page numbers on separate lines
    text = re.sub(r'\b[Pp]age\s*\d+\s*of\s*\d+\b', '', text)  # "Page X of Y"
    text = re.sub(r'\b[Pp]age\s*\d+\b', '', text)  # "Page X"
    
    return text.strip()

def get_metadata_from_path(file_path: str, input_dir: str) -> Dict[str, str]:
    """Extract metadata from file path based on directory structure"""
    rel_path = os.path.relpath(file_path, input_dir)
    parts = Path(rel_path).parts
    
    metadata = {
        "Fichier source": os.path.basename(file_path)
    }
    
    # Extract year, module and type based on directory structure
    if len(parts) >= 1:
        metadata["Année"] = parts[0]
    if len(parts) >= 2:
        metadata["Module"] = parts[1]
    if len(parts) >= 3 and parts[-1] != os.path.basename(file_path):
        metadata["Type"] = parts[2]
    
    return metadata

def create_yaml_header(metadata: Dict[str, str]) -> str:
    """Create YAML header from metadata"""
    yaml_header = "---\n"
    for key, value in metadata.items():
        yaml_header += f"{key}: {value}\n"
    yaml_header += "---\n\n"
    return yaml_header

def process_file(file_path: str, input_dir: str, output_dir: str) -> bool:
    """Process a single file and save the extracted text"""
    try:
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # Extract text based on file extension
        if file_ext in ['.pdf']:
            text = extract_text_from_pdf(file_path)
        elif file_ext in ['.jpg', '.jpeg', '.png', '.tiff', '.tif', '.bmp', '.gif']:
            text = extract_text_from_image(file_path)
        elif file_ext in ['.docx']:
            text = extract_text_from_docx(file_path)
        elif file_ext in ['.doc']:
            text = extract_text_from_doc(file_path)
        elif file_ext in ['.pptx']:
            text = extract_text_from_pptx(file_path)
        elif file_ext in ['.ppt']:
            text = extract_text_from_ppt(file_path)
        else:
            logger.warning(f"Unsupported file type: {file_path}")
            return False
        
        # Clean the extracted text
        cleaned_text = clean_text(text)
        
        if not cleaned_text:
            logger.warning(f"No text extracted from: {file_path}")
            return False
        
        # Get metadata and create YAML header
        metadata = get_metadata_from_path(file_path, input_dir)
        yaml_header = create_yaml_header(metadata)
        
        # Prepare output path (maintaining directory structure)
        rel_path = os.path.relpath(file_path, input_dir)
        output_path = os.path.join(output_dir, os.path.splitext(rel_path)[0] + '.txt')
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # Write text with YAML header to output file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(yaml_header + cleaned_text)
        
        logger.info(f"Successfully processed: {file_path}")
        return True
    
    except Exception as e:
        logger.error(f"Error processing {file_path}: {str(e)}")
        logger.debug(traceback.format_exc())
        return False

def process_directory(input_dir: str, output_dir: str) -> Tuple[int, int]:
    """Process all files in the input directory recursively"""
    success_count = 0
    fail_count = 0
    
    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Supported file extensions
    supported_extensions = [
        '.pdf', '.jpg', '.jpeg', '.png', '.tiff', '.tif', '.bmp', '.gif',
        '.docx', '.doc', '.pptx', '.ppt'
    ]
    
    # Walk through the directory tree
    for root, _, files in os.walk(input_dir):
        for file in files:
            file_path = os.path.join(root, file)
            file_ext = os.path.splitext(file)[1].lower()
            
            if file_ext in supported_extensions:
                if process_file(file_path, input_dir, output_dir):
                    success_count += 1
                else:
                    fail_count += 1
    
    return success_count, fail_count


input_dir = "C:/Users/INFO STOCK 2022/OneDrive/Bureau/Project Medecine/Annales" 
output_dir = "C:/Users/INFO STOCK 2022/OneDrive/Bureau/S6/Textes_Extraits_Annales" 

# Lancer l'extraction directement
logger.info(f"Starting text extraction from {input_dir} to {output_dir}")

# Check if input directory exists
if not os.path.isdir(input_dir):
    logger.error(f"Input directory does not exist: {input_dir}")
    print(f"Le dossier d'entrée n'existe pas: {input_dir}")
else:
    # Process all files
    success_count, fail_count = process_directory(input_dir, output_dir)
    
    logger.info(f"Extraction complete. Successfully processed: {success_count}, Failed: {fail_count}")
    print(f"Extraction terminée. Fichiers traités avec succès: {success_count}, Échecs: {fail_count}")
